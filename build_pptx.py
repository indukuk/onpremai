"""Build an editable PowerPoint from the pitch deck content using python-pptx."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
DARKER_BG = RGBColor(0x0F, 0x0F, 0x23)
MID_BG = RGBColor(0x16, 0x21, 0x3E)
ACCENT = RGBColor(0x00, 0xD4, 0xAA)
WHITE = RGBColor(0xEA, 0xEA, 0xEA)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
TABLE_BG = RGBColor(0x0F, 0x34, 0x60)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, top=Inches(0.5), left=Inches(0.7), width=Inches(12), size=Pt(36)):
    txBox = slide.shapes.add_textbox(left, top, width, Pt(60))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = ACCENT
    return txBox


def add_subtitle(slide, text, top=Inches(1.2), left=Inches(0.7), width=Inches(12), size=Pt(22)):
    txBox = slide.shapes.add_textbox(left, top, width, Pt(40))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = WHITE
    return txBox


def add_body(slide, lines, top=Inches(2.0), left=Inches(0.7), width=Inches(11.5), size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
    return txBox


def add_quote(slide, text, top=Inches(5.5), left=Inches(0.7), width=Inches(11)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = LIGHT_GRAY
    return txBox


def add_table_slide(slide, title, headers, rows, top=Inches(2.0)):
    if title:
        add_title(slide, title)
    num_rows = len(rows) + 1
    num_cols = len(headers)
    left = Inches(0.7)
    width = Inches(11.9)
    row_height = Inches(0.5)
    tbl_height = row_height * num_rows

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, tbl_height)
    table = table_shape.table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        cell.fill.solid()
        cell.fill.fore_color.rgb = MID_BG

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_BG


def arch_box(slide, left, top, width, height, title, items, box_color=MID_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = box_color
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1)

    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), height - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.space_before = Pt(2)


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, DARKER_BG)
add_title(slide, "The Autonomous Compliance Workforce", top=Inches(2.5), size=Pt(44))
add_subtitle(slide, "Every User gets an AI agent that does their compliance work for them", top=Inches(3.5), size=Pt(24))
add_body(slide, [
    "Agents talk to each other. The system gets smarter daily.",
    "Deploys anywhere — including air-gapped.",
], top=Inches(4.5), size=Pt(18))

# ============================================================
# SLIDE 2: The Problem
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, DARK_BG)
add_title(slide, "The Problem")
add_subtitle(slide, "Compliance is a $45B+ market built on human Powered Work")
add_body(slide, [
    "• $3-5M/year per mid-size company — still 80% manual labor",
    "• Control owners lose 5-10 hrs/month on tasks they don't understand",
    "• Compliance managers spend days chasing people, not managing compliance",
    "• Most regulated industries (defense, finance, healthcare) are locked out of modern tooling",
    "  — data can't leave their network",
], top=Inches(2.2))

# ============================================================
# SLIDE 3: Our Solution
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, TABLE_BG)
add_title(slide, "Our Solution")
add_subtitle(slide, "An autonomous AI workforce that runs your compliance program")
add_body(slide, [
    "We don't give companies a better dashboard.",
    "",
    "We give every person an AI agent that does their compliance work",
    "for them — sends messages, collects evidence, evaluates controls,",
    "escalates issues, and coordinates with other agents.",
    "",
    "The Users never needs to understand how our platform works",
], top=Inches(2.5))

# ============================================================
# SLIDE 4: Three Innovations
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, DARK_BG)
add_title(slide, "What Makes This Possible")
add_subtitle(slide, "Three innovations, integrated from day one")
add_table_slide(slide, "", ["#", "Innovation", "What It Means"], [
    ["1", "Shadow AI", "Every user gets a personal agent that acts on their behalf"],
    ["2", "Multi-Agent Coordination", "Agents talk to each other to get work done collaboratively"],
    ["3", "Self-Improving Engine", "3-layer deterministic + LLM pipeline that gets smarter daily"],
], top=Inches(2.2))

# ============================================================
# SLIDE 5: Shadow AI
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, DARK_BG)
add_title(slide, "Shadow AI — The Core Innovation")
add_subtitle(slide, "Every user gets a personal compliance agent. Zero knowledge required.")
add_body(slide, [
    'Agent → Sarah (Slack): "Hey Sarah — I need a Q1 access review',
    '  export from Okta. Want me to ask Mike in IT, or do you have access?"',
    '',
    'Sarah: "Ask Mike"',
    '',
    "Agent → Mike's Agent: [coordinates automatically]",
    '',
    'Mike\'s Agent → Mike (Teams): "Sarah needs an Okta export.',
    '  Takes 2 min — here\'s how."',
    '',
    "Mike: pulls the export → Agent uploads → evaluates → done.",
    '',
    "Sarah never opened the compliance platform. Mike never heard 'SOC 2.'",
    "Total human effort: 2 messages + 2 minutes.",
], top=Inches(2.0), size=Pt(15))

# ============================================================
# SLIDE 6: What the Agent Does
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, DARK_BG)
add_title(slide, "What the Agent Does Behind the Scenes")
add_body(slide, [
    "• Translates compliance requirements into plain language actions",
    "• Knows who has the data and how to get it",
    "• Sends messages, emails, notifications on the user's behalf",
    "• Tracks deadlines and follows up automatically",
    "• Learns each user's communication style and preferences",
    "• Escalates with graduated urgency when things are stuck",
    "• Gets better over time — learns what works for each person",
], top=Inches(1.8))
add_quote(slide, "The agent IS the interface. The platform disappears.", top=Inches(5.5))

# ============================================================
# SLIDE 7: Multi-Agent Network
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, DARK_BG)
add_title(slide, "Multi-Agent Network")
add_subtitle(slide, "Agents don't just assist — they form an autonomous workforce")
add_body(slide, [
    "  [CISO Agent] ◄──► [CompMgr Agent] ◄──► [Owner Agents (Sarah, John)]",
    "                           │                         │",
    "                    [Auditor Agent]          [IT/HR/Dev Agents]",
    "",
    "  Every arrow = agent-to-agent communication",
    "  Humans just approve and answer the occasional question.",
    "",
    "  • 5 agents: basic coordination, some automated evidence flow",
    "  • 20 agents: compliance program runs semi-autonomously",
    "  • 50+ agents: full autonomous workforce, humans manage by exception",
], top=Inches(2.0), size=Pt(16))

# ============================================================
# SLIDE 8: Evaluation Engine
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, DARK_BG)
add_title(slide, "The Evaluation Engine")
add_subtitle(slide, "3-Layer Pipeline: Rules → LLM Judgment → Deterministic Score")
add_table_slide(slide, "", ["Layer", "What", "Cost", "Reproducibility"], [
    ["1", "Deterministic rules (8 types)", "$0", "100%"],
    ["2", "Bounded LLM judgment (ambiguous items only)", "Low", "High"],
    ["3", "Deterministic scoring formula", "$0", "100%"],
], top=Inches(2.2))
add_body(slide, [
    "60-70% of criteria resolve in Layer 1 — zero LLM cost",
    "Layer 2 asks specific questions with rubrics, not 'evaluate this entire control'",
    "Layer 3 always produces the same score given same inputs. Always.",
], top=Inches(4.5), size=Pt(16))

# ============================================================
# SLIDE 9: Why This Matters
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, DARK_BG)
add_title(slide, "Why This Matters")
add_table_slide(slide, "", ["Metric", "Pure LLM (competitors)", "Our 3-Layer Pipeline"], [
    ["Reproducibility", "~70-80%", "97-99%"],
    ["Cost per evaluation", "$$$ (every criteria hits LLM)", "$ (60-70% free)"],
    ["Explainability", '"AI said compliant"', "Rule X passed because Y"],
    ["Speed", "Seconds per criterion", "Milliseconds (rules)"],
    ["Audit trail", "Black box", "Every step traceable"],
], top=Inches(1.8))
add_quote(slide, "Auditors can trust it. CFOs can afford it. Regulators can inspect it.", top=Inches(5.5))

# ============================================================
# SLIDE 10: Self-Improving Observer
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, DARK_BG)
add_title(slide, "Self-Improving System — The Observer")
add_subtitle(slide, "Gets smarter every day. No code deploy needed.")
add_table_slide(slide, "", ["Signal", "Meaning", "Auto-Fix"], [
    ["55% escalation rate", "Model struggling", "Route to stronger tier"],
    ["Confidence trending ↓", "Prompt degrading", "Rewrite via canary test"],
    ["Parse failures >15%", "Output format drift", "Adjust constraints"],
    ["Score variance >15%", "Potential bias", "Alert + investigation"],
], top=Inches(2.2))
add_body(slide, [
    "After 6 months: measurably better than day 1.",
    "Competitors: identical to day 1.",
], top=Inches(5.2), size=Pt(18))

# ============================================================
# SLIDE 11: Graduated Autonomy
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, DARK_BG)
add_title(slide, "Graduated Autonomy — Not a Loose Cannon")
add_table_slide(slide, "", ["Tier", "What", "Safety"], [
    ["1 — Auto-apply", "Routing changes, thresholds", "Confidence ≥ 0.80, 20+ samples"],
    ["2 — Canary first", "Prompt rewrites, model swaps", "20% traffic, 30+ samples, 4+ hrs"],
    ["3 — Human approval", "Model removal, policy changes", "Notify and wait"],
], top=Inches(2.0))
add_body(slide, [
    "Circuit breaker: 3+ rollbacks in 6 hours → all auto-applies stop",
    "The system improves aggressively but fails safely.",
], top=Inches(4.8), size=Pt(18))

# ============================================================
# SLIDE 12: AI Governance
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, DARK_BG)
add_title(slide, "AI Governance Built In")
add_subtitle(slide, "The system audits itself. For SR 11-7, EU AI Act, and customer trust.")
add_body(slide, [
    "Automated model governance (no human maintenance):",
    "• Model inventory — auto-tracked from routing config",
    "• Drift detection — KS test on weekly confidence distributions",
    "• Bias monitoring — cross-tenant score variance flagging",
    "• Performance trending — per model, per task, over time",
    "• Decision audit trail — every observer action: what, why, outcome",
    "• Weekly/monthly governance reports — structured for auditors",
], top=Inches(2.2))
add_quote(slide, 'Auditor: "Show me your AI governance."  Us: GET /observer/governance-report — generated automatically since deployment.', top=Inches(5.8))

# ============================================================
# SLIDE 13: Deploy Anywhere
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, DARK_BG)
add_title(slide, "Deploy Anywhere — First-Class Air-Gapped")
add_table_slide(slide, "", ["Cloud-Only", "Hybrid", "Air-Gapped"], [
    ["All LLMs deployed on AWS", "Local eval, cloud for complex", "Everything on-prem"],
    ["Fastest to start", "Best of both", "Total data control"],
], top=Inches(2.0))
add_body(slide, [
    "Same code. Same Docker images. Different routing.yaml.",
    "",
    "Air-gapped deployment:",
    "  tar xzf compliance-v1.5.0-offline.tar.gz && ./install.sh",
    "  Zero data leaves the building.",
], top=Inches(3.8), size=Pt(16))

# ============================================================
# SLIDE 14: On-Prem Infrastructure Costs
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, DARK_BG)
add_title(slide, "On-Premises Infrastructure Requirements")

add_subtitle(slide, "Application Tier (CPU-only VMs)", top=Inches(1.3), size=Pt(18))

add_table_slide(slide, "", ["VM", "Spec", "Role", "Est. Cost/mo"], [
    ["App VM 1", "8 vCPU, 32 GB RAM, 200 GB SSD", "All 8 services (Docker Compose)", "~$300"],
    ["App VM 2", "8 vCPU, 32 GB RAM, 200 GB SSD", "HA replica / horizontal scale", "~$300"],
    ["DB VM", "8 vCPU, 64 GB RAM, 1 TB NVMe", "PostgreSQL + pgvector, Redis, MinIO", "~$500"],
], top=Inches(1.8))

# LLM tier table
txBox = slide.shapes.add_textbox(Inches(0.7), Inches(3.6), Inches(12), Pt(30))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "LLM Tier (GPU VMs) — separate from app tier"
p.font.size = Pt(18)
p.font.color.rgb = WHITE

num_rows = 4
num_cols = 4
left = Inches(0.7)
width = Inches(11.9)
row_height = Inches(0.5)
tbl_height = row_height * num_rows

table_shape = slide.shapes.add_table(num_rows, num_cols, left, Inches(4.1), width, tbl_height)
table = table_shape.table

headers = ["VM", "Spec", "Models Hosted", "Est. Cost/mo"]
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    cell.fill.solid()
    cell.fill.fore_color.rgb = MID_BG

llm_rows = [
    ["LLM Strong", "4x A100 80GB (or 2x H100)", "DeepSeek-V3 / Qwen-72B (1M context)", "~$8,000-12,000"],
    ["LLM Mid", "1x A100 80GB (or 2x A6000)", "Mistral-22B / Qwen-32B", "~$2,500-4,000"],
    ["LLM Fast", "1x A10 24GB (or RTX 4090)", "Phi-3 / Qwen-7B / Llama-8B", "~$500-1,000"],
]
for r_idx, row in enumerate(llm_rows):
    for c_idx, val in enumerate(row):
        cell = table.cell(r_idx + 1, c_idx)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_BG

add_body(slide, [
    "Why multiple LLMs?",
    "  Strong (1M context): Complex policy analysis, long document reasoning",
    "  Mid: Standard evaluations, evidence review, agent coordination",
    "  Fast: Classification, routing, simple checks (60-70% of all calls)",
    "",
    "Total on-prem: 5-6 VMs, ~$12K-18K/mo (vs. $50K-100K/yr SaaS + data sovereignty risk)",
], top=Inches(6.0), size=Pt(14))

# ============================================================
# SLIDE 15: AWS Architecture
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, DARK_BG)
add_title(slide, "Architecture: AWS Deployment")

arch_box(slide, Inches(0.3), Inches(1.5), Inches(3.0), Inches(2.5),
    "IDENTITY & AUTH", ["Azure AD / Okta (SAML/OIDC)", "AWS Cognito (federated)", "IAM Roles (S2S)", "Secrets Manager"])

arch_box(slide, Inches(3.6), Inches(1.5), Inches(5.5), Inches(2.5),
    "ECS FARGATE (private subnets)", ["llm-gateway → Bedrock (LLMs)", "agent-eval | compliance-assistant",
    "memory-svc | observer", "preprocessor | sandbox-svc"])

arch_box(slide, Inches(9.4), Inches(1.5), Inches(3.6), Inches(2.5),
    "COMMS & CHANNELS", ["Slack API", "Microsoft Teams", "Email (SES)", "Webhooks"])

arch_box(slide, Inches(0.3), Inches(4.3), Inches(3.0), Inches(2.8),
    "CUSTOMER APPS", ["Jira / ServiceNow", "Splunk / Datadog", "Workday HCM", "Salesforce"])

arch_box(slide, Inches(3.6), Inches(4.3), Inches(5.5), Inches(2.8),
    "MANAGED SERVICES", ["RDS PostgreSQL + pgvector", "S3 (evidence artifacts)",
    "ElastiCache Redis", "Cloud Map (discovery)", "VPC + ALB + PrivateLink + VPN", "Bedrock LLM and Agents"])

arch_box(slide, Inches(9.4), Inches(4.3), Inches(3.6), Inches(2.8),
    "FILE STORAGE", ["SharePoint Online", "Google Drive", "Box / Dropbox", "Confluence"])

# ============================================================
# SLIDE 16: On-Prem Architecture
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, DARK_BG)
add_title(slide, "Architecture: Fully On-Premises (Air-Gapped)")

arch_box(slide, Inches(0.3), Inches(1.5), Inches(3.0), Inches(2.5),
    "IDENTITY", ["Active Directory / LDAP", "Keycloak (OIDC/SAML)", "HashiCorp Vault (secrets)", "Internal CA (mTLS)"])

arch_box(slide, Inches(3.6), Inches(1.5), Inches(5.5), Inches(2.5),
    "DOCKER COMPOSE / K8s (app tier)", ["llm-gateway → Ollama / vLLM (GPU)", "agent-eval | compliance-assistant",
    "memory-svc | observer", "preprocessor | sandbox-svc"])

arch_box(slide, Inches(9.4), Inches(1.5), Inches(3.6), Inches(2.5),
    "COMMS (internal only)", ["Teams (on-prem)", "SMTP relay", "Mattermost", "Cisco Webex"])

arch_box(slide, Inches(0.3), Inches(4.3), Inches(3.0), Inches(2.8),
    "CUSTOMER APPS", ["Jira Server", "ServiceNow (on-prem)", "Archer / RSA", "SAP GRC", "Workday HCM"])

arch_box(slide, Inches(3.6), Inches(4.3), Inches(5.5), Inches(2.8),
    "SELF-HOSTED INFRA", ["PostgreSQL + pgvector", "MinIO (S3-compatible)",
    "Redis", "Tesseract (OCR)", "Internal DNS (discovery)", "Reverse proxy (nginx/HAProxy)", "No egress — fully air-gapped"])

arch_box(slide, Inches(9.4), Inches(4.3), Inches(3.6), Inches(2.8),
    "FILE STORAGE", ["NFS / CIFS shares", "SharePoint (on-prem)", "Documentum", "Network drives"])

# ============================================================
# SLIDE 17: Vision
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, DARKER_BG)
add_title(slide, "The Vision", top=Inches(1.5), size=Pt(40))
add_body(slide, [
    "Today (everyone else):",
    "  Humans do compliance work. Tools help them do it faster.",
    "",
    "Tomorrow (us):",
    "  AI agents do compliance work. Humans approve and course-correct.",
    "",
    "",
    "The compliance platform disappears.",
    "The agents are the product.",
    "",
    "Nobody learns a tool. Nobody navigates a dashboard.",
    "Nobody reads a framework document.",
    "The agent handles all of it — and gets better at it every single day.",
], top=Inches(2.5), size=Pt(18))

# ============================================================
# Save
# ============================================================
output_path = "pitch-deck-slides.pptx"
prs.save(output_path)
print(f"Editable PPTX saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
