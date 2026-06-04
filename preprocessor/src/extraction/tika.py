"""Document extraction backend using python-docx and python-pptx.

Handles Word (.docx) and PowerPoint (.pptx) documents natively without
requiring an external Tika server. Uses the python-docx and python-pptx
libraries for direct parsing.
"""

from __future__ import annotations

import io

import structlog

from src.extraction.base import ExtractionBackend, ExtractionResult

logger = structlog.get_logger(__name__)


class TikaBackend(ExtractionBackend):
    """Document extraction for Word and PowerPoint files.

    Despite the name (kept for config compatibility), this backend uses
    python-docx and python-pptx directly rather than Apache Tika server.
    This avoids the need for a separate JVM-based service.
    """

    async def extract_text(self, file_bytes: bytes, **kwargs: object) -> ExtractionResult:
        """Extract text from Word or PowerPoint documents.

        Args:
            file_bytes: Raw document bytes (.docx or .pptx).
            **kwargs: Additional options:
                - file_type (str): "docx" or "pptx" to hint document type.

        Returns:
            ExtractionResult with extracted text.
        """
        file_type = str(kwargs.get("file_type", "docx"))

        if file_type in ("pptx", "powerpoint"):
            return await self._extract_pptx(file_bytes)
        return await self._extract_docx(file_bytes)

    async def _extract_docx(self, file_bytes: bytes) -> ExtractionResult:
        """Extract text from a Word document."""
        try:
            from docx import Document

            doc = Document(io.BytesIO(file_bytes))
            paragraphs: list[str] = []

            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    paragraphs.append(text)

            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        paragraphs.append(row_text)

            full_text = "\n".join(paragraphs)

            logger.info(
                "docx_extraction_complete",
                paragraphs=len(paragraphs),
                text_length=len(full_text),
            )

            return ExtractionResult(
                text=full_text,
                page_count=max(1, len(full_text) // 3000),  # Approximate page count
                confidence=1.0,
                metadata={"backend": "python-docx", "paragraph_count": len(paragraphs)},
            )

        except ImportError:
            return ExtractionResult.failure("python-docx not installed")
        except Exception as exc:
            logger.error(
                "docx_extraction_failed",
                error=str(exc),
                error_type=type(exc).__name__,
            )
            return ExtractionResult.failure(f"Word extraction failed: {exc}")

    async def _extract_pptx(self, file_bytes: bytes) -> ExtractionResult:
        """Extract text from a PowerPoint presentation."""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(file_bytes))
            slides_text: list[str] = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_parts: list[str] = [f"[Slide {slide_num}]"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_parts.append(text)
                slides_text.append("\n".join(slide_parts))

            full_text = "\n\n".join(slides_text)

            logger.info(
                "pptx_extraction_complete",
                slides=len(prs.slides),
                text_length=len(full_text),
            )

            return ExtractionResult(
                text=full_text,
                page_count=len(prs.slides),
                confidence=1.0,
                metadata={"backend": "python-pptx", "slide_count": len(prs.slides)},
            )

        except ImportError:
            return ExtractionResult.failure("python-pptx not installed")
        except Exception as exc:
            logger.error(
                "pptx_extraction_failed",
                error=str(exc),
                error_type=type(exc).__name__,
            )
            return ExtractionResult.failure(f"PowerPoint extraction failed: {exc}")

    def is_available(self) -> bool:
        """Check if python-docx is installed (minimum requirement)."""
        try:
            import docx  # noqa: F401

            return True
        except ImportError:
            return False

    @property
    def name(self) -> str:
        """Backend name."""
        return "tika"
